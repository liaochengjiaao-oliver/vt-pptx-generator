#!/usr/bin/env python3
"""
Generate VT-styled PPTX from Markdown content files.
Uses the official Virginia Tech template layouts.

MD format:
  ---
  title: Presentation Title
  author: Name
  department: Dept info
  date: 2026-06-02
  ---

  # Section Name        → maroon divider slide
  ## Slide Title         → content slide
  - bullet point        → bullet on current slide
  - **bold:** normal    → bold text support
"""

import sys
import os
import re
import yaml
from datetime import date
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "vt_template.pptx")
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "output")
TEMPLATE_URL = "https://brand.vt.edu/content/dam/brand_vt_edu/downloads/ppt/VT_PresentationTemplate_widescreen_PowerPoint.pptx.zip"

MAROON = RGBColor(0x86, 0x1F, 0x41)
DARK = RGBColor(0x3A, 0x3C, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Layout indices in VT template (Master 0)
LAYOUT_TITLE = 2    # white corner logo, maroon corner (slide 10 style)
LAYOUT_DIVIDER = 8  # 1_maroon background (divider)
LAYOUT_CONTENT = 2  # white corner logo, maroon corner (slide 20 style)


def parse_md(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        text = f.read()

    # Split frontmatter
    parts = text.split('---', 2)
    meta = {}
    body = text
    if len(parts) >= 3 and parts[0].strip() == '':
        meta = yaml.safe_load(parts[1]) or {}
        body = parts[2]

    # Parse body into slides
    slides = []
    current = None

    for line in body.split('\n'):
        line_stripped = line.strip()

        if line_stripped.startswith('# ') and not line_stripped.startswith('## '):
            # Section divider
            slides.append({
                'type': 'divider',
                'title': line_stripped[2:].strip()
            })
            current = None

        elif line_stripped.startswith('## '):
            # New content slide
            current = {
                'type': 'content',
                'title': line_stripped[3:].strip(),
                'bullets': []
            }
            slides.append(current)

        elif line_stripped.startswith('- ') and current and current['type'] == 'content':
            current['bullets'].append(line_stripped[2:].strip())

        elif line_stripped.startswith('  - ') and current and current['type'] == 'content':
            current['bullets'].append({'text': line_stripped[4:].strip(), 'level': 1})

    return meta, slides


def parse_bold(text):
    """Parse **bold** markers. Returns [(text, is_bold), ...]."""
    parts = []
    rest = text
    while '**' in rest:
        i = rest.index('**')
        if i > 0:
            parts.append((rest[:i], False))
        rest = rest[i + 2:]
        j = rest.find('**')
        if j == -1:
            parts.append(('**' + rest, False))
            rest = ''
            break
        parts.append((rest[:j], True))
        rest = rest[j + 2:]
    if rest:
        parts.append((rest, False))
    return parts if parts else [(text, False)]


def set_run(run, text, size_pt, bold=False, color=DARK, font='Calibri'):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_title_slide(prs, meta):
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_TITLE]
    slide = prs.slides.add_slide(layout)

    # PH 0: title
    ph0 = slide.placeholders[0]
    ph0.text = ""
    tf = ph0.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, meta.get('title', 'PRESENTATION TITLE'), 40, color=DARK)

    # PH 10: author + department
    ph10 = slide.placeholders[10]
    ph10.text = ""
    tf10 = ph10.text_frame
    tf10.clear()
    p1 = tf10.paragraphs[0]
    run1 = p1.add_run()
    set_run(run1, meta.get('author', ''), 24, color=MAROON)

    if meta.get('department'):
        p2 = tf10.add_paragraph()
        run2 = p2.add_run()
        set_run(run2, meta['department'], 16, color=DARK)

    # PH 11: date
    if 11 in slide.placeholders:
        ph11 = slide.placeholders[11]
        ph11.text = meta.get('date', date.today().strftime('%Y-%m-%d'))

    return slide


def add_divider_slide(prs, title):
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_DIVIDER]
    slide = prs.slides.add_slide(layout)

    ph0 = slide.placeholders[0]
    ph0.text = ""
    tf = ph0.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, title.upper(), 44, color=WHITE)

    return slide


def add_content_slide(prs, title, bullets):
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_CONTENT]
    slide = prs.slides.add_slide(layout)

    # Title
    ph0 = slide.placeholders[0]
    ph0.text = ""
    tf0 = ph0.text_frame
    tf0.clear()
    p0 = tf0.paragraphs[0]
    run0 = p0.add_run()
    set_run(run0, title.upper(), 25, bold=True, color=DARK)

    # Body bullets
    ph10 = slide.placeholders[10]
    ph10.text = ""
    tf = ph10.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(4)

        level = bullet.get('level', 0) if isinstance(bullet, dict) else 0
        text = bullet.get('text', bullet) if isinstance(bullet, dict) else bullet
        p.level = level

        parts = parse_bold(text)
        for part_text, is_bold in parts:
            run = p.add_run()
            size = 22 if level == 0 else 20
            set_run(run, part_text, size, bold=is_bold, color=DARK)

    return slide


def add_thank_you_slide(prs, text="THANK YOU!", subtext="Questions & Discussion"):
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_DIVIDER]
    slide = prs.slides.add_slide(layout)

    ph0 = slide.placeholders[0]
    ph0.text = ""
    tf = ph0.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, 54, color=WHITE)

    if subtext:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
        run2 = p2.add_run()
        set_run(run2, subtext, 24, color=RGBColor(0xE8, 0x77, 0x22))

    return slide


def generate(md_path):
    if not os.path.exists(TEMPLATE_PATH):
        print("Error: VT template not found at", TEMPLATE_PATH)
        print(f"Download it from:\n  {TEMPLATE_URL}")
        print("Extract the .pptx file and place it as 'vt_template.pptx' in this directory.")
        sys.exit(1)

    meta, slides = parse_md(md_path)

    prs = Presentation(TEMPLATE_PATH)

    # Hide all original template slides
    for slide in prs.slides:
        slide._element.set('show', '0')

    # Build presentation
    add_title_slide(prs, meta)

    for sd in slides:
        if sd['type'] == 'divider':
            add_divider_slide(prs, sd['title'])
        elif sd['type'] == 'content':
            add_content_slide(prs, sd['title'], sd.get('bullets', []))

    # Auto thank-you if last slide isn't one
    add_thank_you_slide(prs)

    # Save
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    base = os.path.splitext(os.path.basename(md_path))[0]
    out = os.path.join(OUTPUT_DIR, f"{base}.pptx")
    prs.save(out)
    print(f"Generated: {out}")
    return out


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python3 generate_pptx.py <content.md>")
        sys.exit(1)
    generate(sys.argv[1])
